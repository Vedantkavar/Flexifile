import streamlit as st
import os
import tempfile
import base64
from PIL import Image
import pandas as pd
import io
# import cairosvg  # Remove this line
from svglib.svglib import svg2rlg  # Add this line
from reportlab.graphics import renderPM  # Add this line
from docx import Document
from docx.shared import Inches
import PyPDF2
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches as PptxInches
import csv
import zipfile
import time

st.set_page_config(
    page_title="Flexifile",
    page_icon="♻️",
    layout="wide"
)

# Custom CSS 
def load_css():
    with open("style.css", "r") as f:
        st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)
load_css()

# App title and description
st.markdown("<div class='main-header'>Universal File Converter</div>", unsafe_allow_html=True)
st.markdown("""
<div class='info-box'>
Convert between various file formats across different domains.
</div>
""", unsafe_allow_html=True)

# Define the format domains and their conversions
format_domains = {
    "Document Formats": {
        "Input Formats": [
            "Microsoft Word (.docx, .doc)",
            "OpenDocument (.odt)",
            "Rich Text Format (.rtf)",
            "Plain Text (.txt)",
            "PDF (.pdf)",
            "HTML (.html, .htm)"
        ],
        "Conversions": {
            "Microsoft Word (.docx, .doc)": ["PDF (.pdf)", "Rich Text Format (.rtf)", "Plain Text (.txt)"],
            "OpenDocument (.odt)": ["PDF (.pdf)", "Microsoft Word (.docx, .doc)"],
            "Rich Text Format (.rtf)": ["Microsoft Word (.docx, .doc)", "PDF (.pdf)"],
            "Plain Text (.txt)": ["Microsoft Word (.docx, .doc)", "PDF (.pdf)"],
            "PDF (.pdf)": ["Microsoft Word (.docx, .doc)", "Plain Text (.txt)"],
            "HTML (.html, .htm)": ["Microsoft Word (.docx, .doc)", "PDF (.pdf)"]
        }
    },
    "Presentation Formats": {
        "Input Formats": [
            "Microsoft PowerPoint (.pptx, .ppt)",
            "OpenDocument Presentation (.odp)",
            "PDF (.pdf)"
        ],
        "Conversions": {
            "Microsoft PowerPoint (.pptx, .ppt)": ["PDF (.pdf)", "PNG/JPEG Images (.png, .jpg)"],
            "OpenDocument Presentation (.odp)": ["Microsoft PowerPoint (.pptx, .ppt)", "PDF (.pdf)"],
            "PDF (.pdf)": ["Microsoft PowerPoint (.pptx, .ppt)", "PNG/JPEG Images (.png, .jpg)"]
        }
    },
    "Spreadsheet Formats": {
        "Input Formats": [
            "Microsoft Excel (.xlsx, .xls)",
            "OpenDocument Spreadsheet (.ods)",
            "CSV (.csv)",
            "TSV (.tsv)"
        ],
        "Conversions": {
            "Microsoft Excel (.xlsx, .xls)": ["CSV (.csv)", "TSV (.tsv)", "PDF (.pdf)"],
            "OpenDocument Spreadsheet (.ods)": ["Microsoft Excel (.xlsx, .xls)", "CSV (.csv)"],
            "CSV (.csv)": ["Microsoft Excel (.xlsx, .xls)", "TSV (.tsv)"],
            "TSV (.tsv)": ["Microsoft Excel (.xlsx, .xls)", "CSV (.csv)"]
        }
    },
    "Image Formats": {
        "Input Formats": [
            "PNG (.png)",
            "JPEG (.jpg, .jpeg)",
            "BMP (.bmp)",
            "TIFF (.tiff, .tif)",
            "GIF (.gif)",
            "WebP (.webp)"
        ],
        "Conversions": {
            "PNG (.png)": ["JPEG (.jpg, .jpeg)", "BMP (.bmp)", "TIFF (.tiff, .tif)", "WebP (.webp)"],
            "JPEG (.jpg, .jpeg)": ["PNG (.png)", "BMP (.bmp)", "TIFF (.tiff, .tif)", "WebP (.webp)"],
            "BMP (.bmp)": ["PNG (.png)", "JPEG (.jpg, .jpeg)", "TIFF (.tiff, .tif)"],
            "TIFF (.tiff, .tif)": ["PNG (.png)", "JPEG (.jpg, .jpeg)", "BMP (.bmp)"],
            "GIF (.gif)": ["PNG (.png)", "JPEG (.jpg, .jpeg)"],
            "WebP (.webp)": ["PNG (.png)", "JPEG (.jpg, .jpeg)"]
        }
    },
    "Vector Graphics": {
        "Input Formats": [
            "SVG (.svg)",
            "EPS (.eps)",
            "PDF (.pdf)"
        ],
        "Conversions": {
            "SVG (.svg)": ["PNG (.png)", "JPEG (.jpg, .jpeg)", "PDF (.pdf)"],
            "EPS (.eps)": ["PDF (.pdf)", "PNG (.png)"],
            "PDF (.pdf)": ["PNG (.png)", "JPEG (.jpg, .jpeg)"]
        }
    }
}

# File extension mappings
extension_map = {
    "Microsoft Word (.docx, .doc)": ".docx",
    "OpenDocument (.odt)": ".odt",
    "Rich Text Format (.rtf)": ".rtf",
    "Plain Text (.txt)": ".txt",
    "PDF (.pdf)": ".pdf",
    "HTML (.html, .htm)": ".html",
    "Microsoft PowerPoint (.pptx, .ppt)": ".pptx",
    "OpenDocument Presentation (.odp)": ".odp",
    "PNG/JPEG Images (.png, .jpg)": ".png",
    "Microsoft Excel (.xlsx, .xls)": ".xlsx",
    "OpenDocument Spreadsheet (.ods)": ".ods",
    "CSV (.csv)": ".csv",
    "TSV (.tsv)": ".tsv",
    "PNG (.png)": ".png",
    "JPEG (.jpg, .jpeg)": ".jpg",
    "BMP (.bmp)": ".bmp",
    "TIFF (.tiff, .tif)": ".tiff",
    "GIF (.gif)": ".gif",
    "WebP (.webp)": ".webp",
    "SVG (.svg)": ".svg",
    "EPS (.eps)": ".eps"
}

# Function to get file extension from filename
def get_file_extension(filename):
    return os.path.splitext(filename)[1].lower()

# Function to create a download link
def create_download_link(file_path, filename):
    with open(file_path, "rb") as f:
        bytes_data = f.read()
    b64 = base64.b64encode(bytes_data).decode()
    href = f'<a href="data:application/octet-stream;base64,{b64}" download="{filename}" class="download-btn">Download Converted File</a>'
    return href

# Function to create a download link for in-memory data
def create_download_link_from_bytes(bytes_data, filename):
    b64 = base64.b64encode(bytes_data).decode()
    href = f'<a href="data:application/octet-stream;base64,{b64}" download="{filename}" class="download-btn">Download Converted File</a>'
    return href

# Function to create a download link for multiple files (ZIP)
def create_zip_download_link(file_paths, zip_filename):
    # Create a temporary file to store the ZIP
    with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as temp_zip:
        with zipfile.ZipFile(temp_zip.name, 'w') as zipf:
            for file_path in file_paths:
                zipf.write(file_path, os.path.basename(file_path))
    
    # Create download link for the ZIP file
    return create_download_link(temp_zip.name, zip_filename)

# Conversion functions
def convert_document(input_file, input_format, output_format, input_filename):
    # Create a temporary directory for file operations
    temp_dir = tempfile.mkdtemp()
    
    # Save the uploaded file to the temporary directory
    input_path = os.path.join(temp_dir, input_filename)
    with open(input_path, "wb") as f:
        f.write(input_file.getbuffer())
    
    # Get base filename without extension
    base_filename = os.path.splitext(input_filename)[0]
    
    # Document conversions
    if input_format == "Microsoft Word (.docx, .doc)" and output_format == "PDF (.pdf)":
        # Word to PDF conversion
        from docx2pdf import convert
        output_path = os.path.join(temp_dir, f"{base_filename}.pdf")
        convert(input_path, output_path)
        return output_path, f"{base_filename}.pdf"
    
    elif input_format == "PDF (.pdf)" and output_format == "Microsoft Word (.docx, .doc)":
        # PDF to Word conversion (simplified)
        output_path = os.path.join(temp_dir, f"{base_filename}.docx")
        doc = Document()
        
        # Extract text from PDF
        pdf_file = open(input_path, 'rb')
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text = page.extract_text()
            doc.add_paragraph(text)
        
        doc.save(output_path)
        pdf_file.close()
        return output_path, f"{base_filename}.docx"
    
    elif input_format == "Plain Text (.txt)" and output_format == "Microsoft Word (.docx, .doc)":
        # TXT to Word conversion
        output_path = os.path.join(temp_dir, f"{base_filename}.docx")
        doc = Document()
        
        with open(input_path, 'r', encoding='utf-8') as txt_file:
            text = txt_file.read()
            doc.add_paragraph(text)
        
        doc.save(output_path)
        return output_path, f"{base_filename}.docx"
    
    elif input_format == "Plain Text (.txt)" and output_format == "PDF (.pdf)":
        # TXT to PDF conversion (via Word)
        docx_path = os.path.join(temp_dir, f"{base_filename}.docx")
        output_path = os.path.join(temp_dir, f"{base_filename}.pdf")
        
        # First convert to DOCX
        doc = Document()
        with open(input_path, 'r', encoding='utf-8') as txt_file:
            text = txt_file.read()
            doc.add_paragraph(text)
        doc.save(docx_path)
        
        # Then convert DOCX to PDF
        from docx2pdf import convert
        convert(docx_path, output_path)
        return output_path, f"{base_filename}.pdf"
    
    # Add more document conversions as needed
    
    # Default case - return error
    return None, "Conversion not supported"

def convert_presentation(input_file, input_format, output_format, input_filename):
    # Create a temporary directory for file operations
    temp_dir = tempfile.mkdtemp()
    
    # Save the uploaded file to the temporary directory
    input_path = os.path.join(temp_dir, input_filename)
    with open(input_path, "wb") as f:
        f.write(input_file.getbuffer())
    
    # Get base filename without extension
    base_filename = os.path.splitext(input_filename)[0]
    
    # Presentation conversions
    if input_format == "Microsoft PowerPoint (.pptx, .ppt)" and output_format == "PDF (.pdf)":
        # PowerPoint to PDF conversion
        # Note: Direct conversion from PPTX to PDF in Python is limited
        # This is a simplified approach that creates a basic PDF
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from pptx import Presentation
        
        output_path = os.path.join(temp_dir, f"{base_filename}.pdf")
        prs = Presentation(input_path)
        
        # Create a simple PDF with text content from slides
        c = canvas.Canvas(output_path, pagesize=letter)
        for i, slide in enumerate(prs.slides):
            if i > 0:  # Add a new page for each slide after the first
                c.showPage()
            
            # Add slide number
            c.setFont("Helvetica", 10)
            c.drawString(500, 750, f"Slide {i+1}")
            
            # Add slide title if available
            c.setFont("Helvetica-Bold", 16)
            title_shape = next((shape for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text), None)
            if title_shape:
                c.drawString(72, 720, title_shape.text_frame.text[:50])
            
            # Add text from other shapes
            c.setFont("Helvetica", 12)
            y_position = 680
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = " ".join(run.text for run in paragraph.runs)
                        if text and text != title_shape.text_frame.text if title_shape else True:
                            c.drawString(72, y_position, text[:80])
                            y_position -= 20
                            if y_position < 72:  # Avoid writing off the page
                                break
        
        c.save()
        st.info("Note: This is a basic text-only conversion. For full fidelity PowerPoint to PDF conversion, consider using desktop software like Microsoft PowerPoint or LibreOffice.")
        return output_path, f"{base_filename}.pdf"
    
    elif input_format == "Microsoft PowerPoint (.pptx, .ppt)" and output_format == "PNG/JPEG Images (.png, .jpg)":
        # PowerPoint to Images conversion
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        
        prs = Presentation(input_path)
        
        # Create a directory for the images
        images_dir = os.path.join(temp_dir, "slides")
        os.makedirs(images_dir, exist_ok=True)
        
        # Create simple image representations of slides
        image_paths = []
        for i, slide in enumerate(prs.slides):
            img_path = os.path.join(images_dir, f"slide_{i+1}.png")
            
            # Create a blank image
            img = Image.new('RGB', (1280, 720), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # Try to use a system font
            try:
                font = ImageFont.truetype("Arial", 24)
                small_font = ImageFont.truetype("Arial", 16)
            except IOError:
                font = ImageFont.load_default()
                small_font = ImageFont.load_default()
            
            # Add slide number
            draw.text((1150, 30), f"Slide {i+1}", fill=(0, 0, 0), font=small_font)
            
            # Add slide title if available
            title_shape = next((shape for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text), None)
            if title_shape:
                draw.text((100, 100), title_shape.text_frame.text[:50], fill=(0, 0, 0), font=font)
            
            # Add text from other shapes
            y_position = 200
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = " ".join(run.text for run in paragraph.runs)
                        if text and text != title_shape.text_frame.text if title_shape else True:
                            draw.text((100, y_position), text[:80], fill=(0, 0, 0), font=small_font)
                            y_position += 30
                            if y_position > 650:  # Avoid writing off the image
                                break
            
            img.save(img_path)
            image_paths.append(img_path)
        
        # Create a ZIP file with all images
        zip_path = os.path.join(temp_dir, f"{base_filename}_slides.zip")
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for img_path in image_paths:
                zipf.write(img_path, os.path.basename(img_path))
        
        st.info("Note: This is a basic text-only conversion. For full fidelity PowerPoint to image conversion, consider using desktop software.")
        return zip_path, f"{base_filename}_slides.zip"
    
    # Add more presentation conversions as needed
    
    # Default case - return error
    return None, "Conversion not supported"

def convert_spreadsheet(input_file, input_format, output_format, input_filename):
    # Create a temporary directory for file operations
    temp_dir = tempfile.mkdtemp()
    
    # Save the uploaded file to the temporary directory
    input_path = os.path.join(temp_dir, input_filename)
    with open(input_path, "wb") as f:
        f.write(input_file.getbuffer())
    
    # Get base filename without extension
    base_filename = os.path.splitext(input_filename)[0]
    
    # Spreadsheet conversions
    if input_format == "Microsoft Excel (.xlsx, .xls)" and output_format == "CSV (.csv)":
        # Excel to CSV conversion
        output_path = os.path.join(temp_dir, f"{base_filename}.csv")
        df = pd.read_excel(input_path)
        df.to_csv(output_path, index=False)
        return output_path, f"{base_filename}.csv"
    
    elif input_format == "Microsoft Excel (.xlsx, .xls)" and output_format == "TSV (.tsv)":
        # Excel to TSV conversion
        output_path = os.path.join(temp_dir, f"{base_filename}.tsv")
        df = pd.read_excel(input_path)
        df.to_csv(output_path, sep='\t', index=False)
        return output_path, f"{base_filename}.tsv"
    
    elif input_format == "CSV (.csv)" and output_format == "Microsoft Excel (.xlsx, .xls)":
        # CSV to Excel conversion
        output_path = os.path.join(temp_dir, f"{base_filename}.xlsx")
        df = pd.read_csv(input_path)
        df.to_excel(output_path, index=False)
        return output_path, f"{base_filename}.xlsx"
    
    elif input_format == "TSV (.tsv)" and output_format == "Microsoft Excel (.xlsx, .xls)":
        # TSV to Excel conversion
        output_path = os.path.join(temp_dir, f"{base_filename}.xlsx")
        df = pd.read_csv(input_path, sep='\t')
        df.to_excel(output_path, index=False)
        return output_path, f"{base_filename}.xlsx"
    
    elif input_format == "CSV (.csv)" and output_format == "TSV (.tsv)":
        # CSV to TSV conversion
        output_path = os.path.join(temp_dir, f"{base_filename}.tsv")
        df = pd.read_csv(input_path)
        df.to_csv(output_path, sep='\t', index=False)
        return output_path, f"{base_filename}.tsv"
    
    elif input_format == "TSV (.tsv)" and output_format == "CSV (.csv)":
        # TSV to CSV conversion
        output_path = os.path.join(temp_dir, f"{base_filename}.csv")
        df = pd.read_csv(input_path, sep='\t')
        df.to_csv(output_path, index=False)
        return output_path, f"{base_filename}.csv"
    
    # Add more spreadsheet conversions as needed
    
    # Default case - return error
    return None, "Conversion not supported"

def convert_image(input_file, input_format, output_format, input_filename):
    # Create a temporary directory for file operations
    temp_dir = tempfile.mkdtemp()
    
    # Save the uploaded file to the temporary directory
    input_path = os.path.join(temp_dir, input_filename)
    with open(input_path, "wb") as f:
        f.write(input_file.getbuffer())
    
    # Get base filename without extension
    base_filename = os.path.splitext(input_filename)[0]
    
    # Image conversions
    try:
        img = Image.open(input_path)
        
        if output_format == "PNG (.png)":
            output_path = os.path.join(temp_dir, f"{base_filename}.png")
            img = img.convert("RGBA")
            img.save(output_path, "PNG")
            return output_path, f"{base_filename}.png"
        
        elif output_format == "JPEG (.jpg, .jpeg)":
            output_path = os.path.join(temp_dir, f"{base_filename}.jpg")
            img = img.convert("RGB")  # Remove alpha for JPG
            img.save(output_path, "JPEG", quality=95)
            return output_path, f"{base_filename}.jpg"
        
        elif output_format == "BMP (.bmp)":
            output_path = os.path.join(temp_dir, f"{base_filename}.bmp")
            img.save(output_path, "BMP")
            return output_path, f"{base_filename}.bmp"
        
        elif output_format == "TIFF (.tiff, .tif)":
            output_path = os.path.join(temp_dir, f"{base_filename}.tiff")
            img.save(output_path, "TIFF")
            return output_path, f"{base_filename}.tiff"
        
        elif output_format == "WebP (.webp)":
            output_path = os.path.join(temp_dir, f"{base_filename}.webp")
            img.save(output_path, "WebP", quality=95)
            return output_path, f"{base_filename}.webp"
        
        # Add more image conversions as needed
    
    except Exception as e:
        st.error(f"Error converting image: {str(e)}")
        return None, f"Error: {str(e)}"
    
    # Default case - return error
    return None, "Conversion not supported"

def convert_vector(input_file, input_format, output_format, input_filename):
    # Create a temporary directory for file operations
    temp_dir = tempfile.mkdtemp()
    
    # Save the uploaded file to the temporary directory
    input_path = os.path.join(temp_dir, input_filename)
    with open(input_path, "wb") as f:
        f.write(input_file.getbuffer())
    
    # Get base filename without extension
    base_filename = os.path.splitext(input_filename)[0]
    
    # Vector graphics conversions
    if input_format == "SVG (.svg)" and output_format == "PNG (.png)":
        # SVG to PNG conversion using svglib
        output_path = os.path.join(temp_dir, f"{base_filename}.png")
        drawing = svg2rlg(input_path)
        renderPM.drawToFile(drawing, output_path, fmt="PNG")
        return output_path, f"{base_filename}.png"
    
    elif input_format == "SVG (.svg)" and output_format == "JPEG (.jpg, .jpeg)":
        # SVG to JPEG conversion using svglib
        output_path = os.path.join(temp_dir, f"{base_filename}.jpg")
        drawing = svg2rlg(input_path)
        renderPM.drawToFile(drawing, output_path, fmt="JPEG")
        return output_path, f"{base_filename}.jpg"
    
    elif input_format == "SVG (.svg)" and output_format == "PDF (.pdf)":
        # SVG to PDF conversion
        # For PDF conversion, you can use reportlab
        from reportlab.graphics import renderPDF
        output_path = os.path.join(temp_dir, f"{base_filename}.pdf")
        drawing = svg2rlg(input_path)
        renderPDF.drawToFile(drawing, output_path)
        return output_path, f"{base_filename}.pdf"
    
    # Add more vector graphics conversions as needed
    
    # Default case - return error
    return None, "Conversion not supported"

# Main app layout

st.markdown("<div class='sub-header'>Step 1: Select File Domain</div>", unsafe_allow_html=True)
domain = st.selectbox("Select File Domain", list(format_domains.keys()))

col1, col2 = st.columns(2)

with col1:
    st.markdown("<div class='sub-header'>Step 2: Select Input Format</div>", unsafe_allow_html=True)
    input_format = st.selectbox("Select Input Format", format_domains[domain]["Input Formats"])

with col2:
    st.markdown("<div class='sub-header'>Step 3: Select Output Format</div>", unsafe_allow_html=True)
    output_format = st.selectbox("Select Output Format", format_domains[domain]["Conversions"].get(input_format, ["No compatible formats"]))

st.markdown("<div class='sub-header'>Step 4: Upload File</div>", unsafe_allow_html=True)
uploaded_file = st.file_uploader("Upload your file", type=None)

# Conversion process
if uploaded_file is not None:
    st.markdown("<div class='sub-header'>Step 5: Convert File</div>", unsafe_allow_html=True)
    
    if st.button("Convert File"):
        with st.spinner("Converting file..."):
            # Determine which conversion function to use based on the domain
            if domain == "Document Formats":
                output_path, output_filename = convert_document(uploaded_file, input_format, output_format, uploaded_file.name)
            elif domain == "Presentation Formats":
                output_path, output_filename = convert_presentation(uploaded_file, input_format, output_format, uploaded_file.name)
            elif domain == "Spreadsheet Formats":
                output_path, output_filename = convert_spreadsheet(uploaded_file, input_format, output_format, uploaded_file.name)
            elif domain == "Image Formats":
                output_path, output_filename = convert_image(uploaded_file, input_format, output_format, uploaded_file.name)
            elif domain == "Vector Graphics":
                output_path, output_filename = convert_vector(uploaded_file, input_format, output_format, uploaded_file.name)
            
            # Add a small delay to simulate processing
            time.sleep(1)
            
            # Check if conversion was successful
            if output_path:
                st.markdown(f"""
                <div class='success-message'>
                    <h3>Conversion Complete!</h3>
                    <p>Your file has been successfully converted from {input_format} to {output_format}.</p>
                </div>
                """, unsafe_allow_html=True)
                
                # Create download link
                st.markdown(create_download_link(output_path, output_filename), unsafe_allow_html=True)
                
                # Preview if possible
                if domain == "Image Formats" or (domain == "Vector Graphics" and "PNG" in output_format or "JPEG" in output_format):
                    st.image(output_path, caption="Preview of converted image", use_column_width=True)
            else:
                st.error(f"Conversion failed: {output_filename}")
                st.info("This conversion may not be fully implemented in this demo or requires additional libraries.")

# Footer with information
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #666;">
    <p>Universal File Converter</p>
</div>
""", unsafe_allow_html=True)

# Add a sidebar with additional information
with st.sidebar:
    st.title("About")
    st.info("""
    This Universal File Converter allows you to convert files between various formats.
    
    **Supported Domains:**
    - Document Formats
    - Presentation Formats
    - Spreadsheet Formats
    - Image Formats
    - Vector Graphics
    
    **Note:** Some conversions may not be supported.
    """)
    
    st.markdown("### How to Use")
    st.markdown("""
    1. Select the file domain
    2. Choose input format
    3. Select output format
    4. Upload your file
    5. Click Convert
    6. Download the converted file
    """)

